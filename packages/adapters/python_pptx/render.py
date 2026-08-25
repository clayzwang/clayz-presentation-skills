#!/usr/bin/env python3
# SPDX-FileCopyrightText: 2026 clayz
# SPDX-License-Identifier: Apache-2.0
"""Render a Clayz render manifest to editable PPTX with python-pptx."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Mapping


def _imports() -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for this runtime route") from exc
    return locals()


def _color(value: Any, api: Mapping[str, Any]) -> Any:
    cleaned = str(value or "000000").lstrip("#")
    if len(cleaned) != 6:
        raise ValueError(f"invalid RGB color: {value!r}")
    return api["RGBColor"].from_string(cleaned.upper())


def _box(options: Mapping[str, Any], api: Mapping[str, Any]) -> tuple[Any, Any, Any, Any]:
    try:
        return tuple(api["Inches"](float(options[key])) for key in ("x", "y", "w", "h"))  # type: ignore[return-value]
    except (KeyError, TypeError, ValueError) as exc:
        raise ValueError("every object requires numeric x, y, w, h options") from exc


def _name(shape: Any, spec: Mapping[str, Any]) -> None:
    copy_id = spec.get("copy_id")
    shape.name = f"COPY::{copy_id}::{spec['object_id']}" if copy_id else str(spec["object_id"])


def _fill_and_line(shape: Any, options: Mapping[str, Any], api: Mapping[str, Any]) -> None:
    fill = options.get("fill")
    if isinstance(fill, Mapping) and fill.get("color"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _color(fill["color"], api)
        if isinstance(fill.get("transparency"), (int, float)):
            shape.fill.transparency = float(fill["transparency"]) / 100.0
    elif fill is None and hasattr(shape, "fill"):
        shape.fill.background()
    line = options.get("line")
    if hasattr(shape, "line"):
        if isinstance(line, Mapping) and line.get("color"):
            shape.line.color.rgb = _color(line["color"], api)
            if isinstance(line.get("width"), (int, float)):
                shape.line.width = api["Pt"](float(line["width"]))
        elif line is None:
            shape.line.fill.background()


def _style_text_frame(text_frame: Any, text: str, options: Mapping[str, Any], api: Mapping[str, Any]) -> None:
    text_frame.clear()
    text_frame.word_wrap = options.get("breakLine", True) is not False
    margin = float(options.get("margin", 0))
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(text_frame, side, api["Inches"](margin))
    valign = str(options.get("valign", "top")).lower()
    text_frame.vertical_anchor = {
        "top": api["MSO_ANCHOR"].TOP,
        "mid": api["MSO_ANCHOR"].MIDDLE,
        "middle": api["MSO_ANCHOR"].MIDDLE,
        "bottom": api["MSO_ANCHOR"].BOTTOM,
    }.get(valign, api["MSO_ANCHOR"].TOP)
    lines = text.split("\n") or [""]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.alignment = {
            "left": api["PP_ALIGN"].LEFT,
            "center": api["PP_ALIGN"].CENTER,
            "right": api["PP_ALIGN"].RIGHT,
            "justify": api["PP_ALIGN"].JUSTIFY,
        }.get(str(options.get("align", "left")).lower(), api["PP_ALIGN"].LEFT)
        if paragraph.runs:
            font = paragraph.runs[0].font
            if options.get("fontFace"):
                font.name = str(options["fontFace"])
            if isinstance(options.get("fontSize"), (int, float)):
                font.size = api["Pt"](float(options["fontSize"]))
            font.bold = bool(options.get("bold", False))
            font.italic = bool(options.get("italic", False))
            if options.get("color"):
                font.color.rgb = _color(options["color"], api)


def _shape_type(name: str, api: Mapping[str, Any]) -> Any:
    normalized = name.replace("-", "").replace("_", "").lower()
    mapping = {
        "rect": api["MSO_AUTO_SHAPE_TYPE"].RECTANGLE,
        "rectangle": api["MSO_AUTO_SHAPE_TYPE"].RECTANGLE,
        "roundrect": api["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
        "ellipse": api["MSO_AUTO_SHAPE_TYPE"].OVAL,
        "oval": api["MSO_AUTO_SHAPE_TYPE"].OVAL,
        "chevron": api["MSO_AUTO_SHAPE_TYPE"].CHEVRON,
        "hexagon": api["MSO_AUTO_SHAPE_TYPE"].HEXAGON,
        "diamond": api["MSO_AUTO_SHAPE_TYPE"].DIAMOND,
        "triangle": api["MSO_AUTO_SHAPE_TYPE"].ISOSCELES_TRIANGLE,
        "arc": api["MSO_AUTO_SHAPE_TYPE"].ARC,
    }
    if normalized not in mapping:
        raise ValueError(f"unsupported shape: {name}")
    return mapping[normalized]


def _chart_type(name: str, api: Mapping[str, Any]) -> Any:
    chart = api["XL_CHART_TYPE"]
    mapping = {
        "bar": chart.BAR_CLUSTERED,
        "column": chart.COLUMN_CLUSTERED,
        "line": chart.LINE_MARKERS,
        "area": chart.AREA,
        "pie": chart.PIE,
        "doughnut": chart.DOUGHNUT,
    }
    if name not in mapping:
        raise ValueError(f"unsupported chart_type: {name}")
    return mapping[name]


def _add_object(slide: Any, spec: Mapping[str, Any], manifest_dir: Path, api: Mapping[str, Any]) -> None:
    options = spec.get("options", {})
    if not isinstance(options, Mapping):
        raise ValueError(f"{spec.get('object_id')}: options must be an object")
    x, y, w, h = _box(options, api)
    kind = spec.get("type")
    if kind == "text":
        shape = slide.shapes.add_textbox(x, y, w, h)
        _style_text_frame(shape.text_frame, str(spec.get("text", "")), options, api)
    elif kind == "shape":
        shape = slide.shapes.add_shape(_shape_type(str(spec.get("shape", "rect")), api), x, y, w, h)
        _fill_and_line(shape, options, api)
    elif kind == "line":
        shape = slide.shapes.add_connector(1, x, y, x + w, y + h)
        _fill_and_line(shape, options, api)
    elif kind == "image":
        raw = spec.get("path")
        if not raw:
            raise ValueError(f"{spec.get('object_id')}: image path is required")
        path = Path(str(raw))
        path = path if path.is_absolute() else manifest_dir / path
        if not path.is_file():
            raise ValueError(f"{spec.get('object_id')}: image not found: {path}")
        shape = slide.shapes.add_picture(str(path), x, y, w, h)
    elif kind == "svg":
        raise ValueError(f"{spec.get('object_id')}: SVG requires an optional converter or a different authoring route")
    elif kind == "table":
        rows = spec.get("rows")
        if not isinstance(rows, list) or not rows or not all(isinstance(row, list) for row in rows):
            raise ValueError(f"{spec.get('object_id')}: rows must be a non-empty matrix")
        cols = max(len(row) for row in rows)
        shape = slide.shapes.add_table(len(rows), cols, x, y, w, h)
        table = shape.table
        for row_index, row in enumerate(rows):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                _style_text_frame(cell.text_frame, str(value), options, api)
    elif kind == "chart":
        series = spec.get("series")
        if not isinstance(series, list) or not series:
            raise ValueError(f"{spec.get('object_id')}: series are required")
        data = api["CategoryChartData"]()
        labels = series[0].get("labels", [])
        data.categories = labels
        for item in series:
            data.add_series(str(item.get("name", "Series")), list(item.get("values", [])))
        shape = slide.shapes.add_chart(_chart_type(str(spec.get("chart_type")), api), x, y, w, h, data)
        chart = shape.chart
        chart.has_legend = bool(options.get("showLegend", len(series) > 1))
        chart.has_title = bool(options.get("showTitle", False))
        if chart.has_legend and options.get("legendPos"):
            pass
        if bool(options.get("showValue", False)):
            for chart_series in chart.series:
                chart_series.has_data_labels = True
    else:
        raise ValueError(f"{spec.get('object_id')}: unsupported object type {kind}")
    _name(shape, spec)


def render(manifest_path: Path, output_path: Path) -> None:
    api = _imports()
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if manifest.get("contract") != "io.clayz.presentation.render-manifest/1.0":
        raise ValueError("unsupported render-manifest contract")
    slides = manifest.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("at least one slide is required")
    presentation = manifest.get("presentation", {})
    pptx = api["Presentation"]()
    if presentation.get("layout") == "LAYOUT_WIDE":
        pptx.slide_width = api["Inches"](13.333333)
        pptx.slide_height = api["Inches"](7.5)
    props = pptx.core_properties
    for attr in ("author", "company", "subject", "title"):
        value = presentation.get(attr)
        if value and hasattr(props, attr):
            setattr(props, attr, str(value))
    seen_slides: set[str] = set()
    seen_objects: set[str] = set()
    for slide_spec in slides:
        slide_id = str(slide_spec.get("slide_id", ""))
        if not slide_id or slide_id in seen_slides:
            raise ValueError("slide_id values must be unique and non-empty")
        seen_slides.add(slide_id)
        slide = pptx.slides.add_slide(pptx.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = _color(slide_spec.get("background", "FFFFFF"), api)
        for spec in slide_spec.get("objects", []):
            object_id = str(spec.get("object_id", ""))
            if not object_id or object_id in seen_objects:
                raise ValueError("object_id values must be globally unique and non-empty")
            seen_objects.add(object_id)
            _add_object(slide, spec, manifest_path.parent, api)
        notes = slide_spec.get("speaker_notes")
        if isinstance(notes, list) and notes:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is not None:
                notes_frame.text = "\n".join(str(item) for item in notes)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pptx.save(output_path)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("manifest", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    try:
        render(args.manifest.resolve(), args.output.resolve())
    except (OSError, ValueError, RuntimeError, json.JSONDecodeError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2
    print(f"wrote {args.output.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
